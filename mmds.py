# TODO: Use proper logger for CodeTimer and other diagnostic output
# TODO: update documentation for gen_mmd_seed_assignment() (and many other funcs) as we are now passing in an mmd config
# TODO: encapsulate Partition and mmd_config into the same data structure
# TODO: simplify recom to maybe one function and simplify complement_or_not logic; it is very confusing
# TODO: fix find_cut() because the tuple return format is very confusing
# TODO: plot vote-seat share curve (search up Shen Github) proportionality
# TODO: rename this file to main.py
# TODO: look into caching the gen_ensemble() function (perhaps with @cache decorator)
# TODO: use linter 
# TODO: use a testing framework to test parts of program (such as running an election)


import random
from disjoint_set import DisjointSet
from linetimer import CodeTimer
from functools import partial
import matplotlib.pyplot as plt
import networkx as nx
from geopandas import GeoSeries
from gerrychain import Graph, MarkovChain, Partition
from gerrychain.accept import always_accept
from gerrychain.updaters import Tally, cut_edges
from matplotlib.colors import ListedColormap
from pptx import Presentation
import itertools
flatten = itertools.chain.from_iterable
from election import Party, Candidate, Ballot, gen_candidates, multi_seat_ranked_choice_election
from voting_models import party_line_voting
from matplotlib.pyplot import hist


SHAPE_FILE = "./state_data/PA/PA.shp"
GRAPH_FILE = "./PAjson_wo_geometry"
PRES_FILE = "./pres_output/mmd_recom3"
SMD_ASSIGNMENT_COL = "CD_2011"
PLOT_INTERVAL = 1 
POP_COL = "TOTPOP"
COLORS = ListedColormap(['#e6194b', '#3cb44b', '#ffe119', '#4363d8', '#f58231',
'#911eb4', '#46f0f0', '#f032e6', '#bcf60c', '#fabebe', '#008080', '#e6beff',
'#9a6324', '#fffac8', '#800000', '#aaffc3', '#808000', '#ffd8b1', '#000075',
'#808080'])


def add_plot_to_pres(prs):
    plt.savefig("/tmp/graph.png")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture("/tmp/graph.png", 0, 0)
    print("slide added")


def plot_partition(partition: Partition, precinct_geometries: GeoSeries, prs, district_reps: dict[int, int], show=False, cmap=COLORS):
    # plt.title("Markov chain step: %d" % chain_step)
    partition.plot(precinct_geometries, cmap=cmap)
    centroids: dict[int, tuple] = get_district_centroids(partition, precinct_geometries)
    for districtID, coord in centroids.items():
        pop_frac = float(partition["population"][districtID]/sum(partition["population"].values())) * sum(district_reps.values())
        plt.text(coord[0], coord[1], "District %d\nPopulation: %d\nPop Frac: %3f/18\n, num reps: %d" % (districtID, partition["population"][districtID], pop_frac, district_reps[districtID]))
    if prs is not None:
        add_plot_to_pres(prs)
    if show:
        plt.show()


def count_components(G: nx.Graph) -> int:
    return len([c for c in nx.connected_components(G)])


def find_cut(graph, tree: Graph, curr_node: int, parent: int, graph_pop: int, pop_rng: tuple) -> tuple:
    sum = graph.nodes[curr_node][POP_COL] 
    for child in tree.neighbors(curr_node):
        if child != parent:
            child_sum, edge_found = find_cut(graph, tree, child, curr_node, graph_pop, pop_rng) 
            if edge_found:
                return (child_sum, edge_found)
            sum += child_sum
    if pop_rng[0] <= sum <= pop_rng[1]:
        return (True, (curr_node, parent))
    elif pop_rng[0] <= graph_pop-sum <= pop_rng[1]:
        return (False, (curr_node, parent))
    return (sum, None)


def split_graph_by_pop(graph: nx.Graph, pop_target: int, graph_pop: int, epsilon: float, node_repeats: int = 30) -> tuple[list[int]]:
    pop_rng = (pop_target * (1 - epsilon), pop_target * (1 + epsilon))
    graph_edges = list(graph.edges)
    for i in range(node_repeats):
        # with CodeTimer("create spanning tree"):
        spanning_tree = rand_spanning_tree(graph, graph_edges)
        root = random.choice(list(spanning_tree.nodes))
        complement_or_not, cut_edge = find_cut(graph, spanning_tree, root, None, graph_pop, pop_rng)
        # print("complement or not = " + str(complement_or_not))
        if cut_edge:
            # print("finished recom after %d random spanning trees on %d nodes" % (i+1, len(spanning_tree.nodes)))
            spanning_tree.remove_edge(cut_edge[0], cut_edge[1])
            comp_1 = [node for node in nx.dfs_postorder_nodes(spanning_tree, source=cut_edge[0])]
            comp_2 = [node for node in nx.dfs_postorder_nodes(spanning_tree, source=cut_edge[1])]
            return((comp_1, comp_2), complement_or_not)
    raise Exception("partitioning failed")


def rand_spanning_tree(graph: nx.Graph, edges):
    spanning_forest = nx.Graph() # can we remove next line by putting graph.nodes into this constructor?
    spanning_forest.add_nodes_from(graph.nodes) 
    random.shuffle(edges)
    ds = DisjointSet()
    for node in spanning_forest.nodes:
        ds.find(node)
    for u, v in edges:
        if not ds.connected(u, v):
            spanning_forest.add_edge(u, v)
            ds.union(u, v)
    return spanning_forest

 
def mmd_recom(partition: Partition, mmd_config: dict[int, int], epsilon: float) -> Partition:
    """
    Computes next partition after one step of MMD ReCom. This is a proposal
    function that can be provided to the gerrychain MarkovChain constructor.
    ReCom algorithm works by merging two adjacent MMD subgraphs, then splitting
    the merged district along a boundary in which each part has a population
    ratio that matches (within the input epsilon's error) the number of
    representatives in that district.

    i.e. if districts A and B are chosen to be merged, and have 3 and 5
    representatives respectively, re-split them into two parts where one part
    has 3/8ths of the combined population and the other part has 5/8ths of the
    combined population.

    To improve intuitiveness of the algorithm, districts will retain the same
    number of representatives throughout all steps of ReCom as specified by
    mmd_config.

    Arguments:
        partition: an MMD partition
        district_reps: dictionary mapping MMD ID to its assigned number of
        representatives
        epsilon: acceptable population error threshold for split
    Returns:
        a new MMD partition
    """

    edge = random.choice(list(partition["cut_edges"]))
    partIDs = (partition.assignment[edge[0]], partition.assignment[edge[1]])
    print("doing recom on districts %d, %d" % (partIDs[0], partIDs[1]))
    merged_subgraph = partition.graph.subgraph(partition.parts[partIDs[0]] | partition.parts[partIDs[1]])
    subgraph_pop = partition["population"][partIDs[0]] + partition["population"][partIDs[1]] 
    subgraph_reps = mmd_config[partIDs[0]] + mmd_config[partIDs[1]]
    pop_target = (float(mmd_config[partIDs[0]])/subgraph_reps)*subgraph_pop
    components, complement_or_not = split_graph_by_pop(merged_subgraph.graph, pop_target, subgraph_pop, epsilon)
    if complement_or_not: # component 0 matches the target pop
        flips = dict.fromkeys(components[0], partIDs[0]) | dict.fromkeys(components[1], partIDs[1]) 
    else:
        flips = dict.fromkeys(components[0], partIDs[1]) | dict.fromkeys(components[1], partIDs[0]) 
    return partition.flip(flips)


def gen_mmd_configs(n_reps: int) -> list[tuple]:
    """
    Generates a list of tuples that represent all possible MMD configurations
    given a number of representatives. Each tuple has 3 elements, where the 0th,
    1st and 2nd indexed elements represent the number of districts with 3, 4,
    and 5 representatives, respectively. Each such unique tuple will hereon be
    referred to as an "MMD config."

    Arguments:
        n_reps: number of representatives in the state
    Returns:
        List of MMD config 3-tuples
    """

    configs = []
    for x1 in range(n_reps//3+1):
        for x2 in range(n_reps//4+1):
            for x3 in range(n_reps//5+1):
                if x1*3 + x2*4 + x3*5 == n_reps:
                    configs.append(dict.fromkeys(range(1, x1+1), 3) 
                        | dict.fromkeys(range(x1+1, x1+x2+1), 4) 
                        | dict.fromkeys(range(x1+x2+1, x1+x2+x3+1), 5))
    return configs


def get_district_centroids(partition: Partition, precinct_geometries: GeoSeries) -> dict[int, tuple]:
    district_centroids = {}
    for part, nodes in partition.parts.items():
        part_geometries = precinct_geometries.loc[precinct_geometries.index[list(nodes)]]
        district_centroids[part] = tuple(part_geometries.unary_union.centroid.coords)[0]
    return district_centroids


def gen_smd_adjacency_graph(partition: Partition) -> nx.Graph:
    """
    Constructs an SMD adjacency graph based on the SMD assignments of the input
    partition.

    Arguments:
        partition: gerrychain Partition initialized with an SMD assignment
    Returns:
        networkx Graph of SMDs
    """

    edgelist = []
    for edge in partition["cut_edges"]:
        edgelist.append((partition.assignment[edge[0]], partition.assignment[edge[1]]))
    return nx.Graph(edgelist)


def cut_smd_adjacency_graph(graph: nx.Graph, mmd_config: dict[int, int], cut_iterations: int = 100000, node_repeats: int = 20) -> nx.Graph:
    """
    Attempts to partition input SMD adjacency graph into connected subgraphs
    with sizes that correspond to each district in the MMD config.  Generates
    random spanning tree of SMD graph and randomly cuts away the # of districts
    in the MMD config - 1. If the number of nodes in each of the cut spanning
    tree's subgraphs match the proportions of the MMD config, it is returned.

    Arguments:
        graph: SMD adjacency graph
        mmd_config: MMD config
        cut_iterations: max number of cut attempts for the randomly generated
        spanning tree
        node_repeats: max number of spanning trees to build
    Returns:
        networkx Graph spanning tree forest that matches input MMD config's
        proportions
    """

    config_sizes = sorted(list(mmd_config.values()))
    for _ in range(node_repeats):
        spanning_tree: nx.Graph = rand_spanning_tree(graph, list(graph.edges))
        edges = list(spanning_tree.edges)
        for _ in range(cut_iterations):
            random_edges: list[tuple] = random.sample(edges, len(mmd_config)-1)
            spanning_tree.remove_edges_from(random_edges)
            forest_sizes = [len(component) for component in nx.connected_components(spanning_tree)]
            if sorted(forest_sizes) == config_sizes:
                return spanning_tree
            spanning_tree.add_edges_from(random_edges)
    raise Exception("partitioning failed after %d random cut iterations after %d node repeats" % (cut_iterations, node_repeats))


def gen_mmd_seed_assignment(smd_partition: Partition, mmd_config: dict[int, int]) -> tuple[dict[int, int], dict[int, int]]:
    """
    Generates the seed MMDs by gluing together adjacent SMD districts such that
    the MMDs that have a population proportion that matches the number of
    representatives are assigned (as specified by mmd_config)

    Generates initial MMD assignment dict that maps precinct IDs to district
    IDs. This will be used as the assignment parameter to the gerrychain
    Partition constructor. Algorithm steps are as follows:

    1. It picks an MMD config
    2. Generates an SMD adjacency graph based on the input partition
    3. Cuts this SMD graph into subgraphs that follow the proportions of the
    selected MMD config
    4. Assigns each precinct in the input partition based on the partitioned SMD
    graph

    Arguments:
        partition: gerrychain Partition initialized with an SMD assignment
    Returns:
        1. assignment dictionary mapping precinct IDs to district IDs
        2. dictionary mapping district IDs to their number of representatives
    """

    smd_graph: nx.Graph = gen_smd_adjacency_graph(smd_partition)
    with CodeTimer("cutting smd graph into proportions specified by config"):
        smd_forest = cut_smd_adjacency_graph(smd_graph, mmd_config)
    components = [c for c in nx.connected_components(smd_forest)]
    prec_assignment = {}
    for districtID, n_reps in mmd_config.items():
        for component in components:
            if len(component) == n_reps:
                component_precIDs = flatten([smd_partition.parts[n] for n in component])
                prec_assignment = prec_assignment | dict.fromkeys(component_precIDs, districtID)
                components.remove(component)
                break
    return prec_assignment


def gen_random_map(chain: MarkovChain) -> Partition:
    for partition in chain:
        continue
    return partition


def gen_ensemble(chain: MarkovChain, num_maps) -> list[Partition]:
    return [gen_random_map(chain) for _ in range(num_maps)]


def run_mmd_election(partition: Partition, districtID: int, mmd_config: dict[int, list[int]]) -> list[Candidate]:
    district_votes: list[Ballot] = []
    with CodeTimer("running gen_candidates"):
        district_candidates: list[Candidate] = gen_candidates(mmd_config[districtID])
    for precID in partition.subgraphs[districtID].nodes:
        with CodeTimer("running party_line_voting"):
            precinct_votes = party_line_voting(partition.graph.nodes[precID], district_candidates)
        district_votes += precinct_votes
    with CodeTimer("running tabulation"):
        return multi_seat_ranked_choice_election(district_votes, district_candidates, mmd_config[districtID])


def run_statewide_election(partition: Partition, mmd_config: dict[int, list[int]]) -> list[Candidate]:
    winners = []
    for districtID in partition.parts.keys():
        with CodeTimer("running mmd election"):
            winners += run_mmd_election(partition, districtID, mmd_config)
    return winners


def run_many_statewide_elections(chain: MarkovChain, mmd_config: dict[int, int], n_elections: int) -> list[list[Candidate]]:
    elections_results: list[list[Candidate]] = []
    with CodeTimer("generating ensemble of size %d" % n_elections):
        ensemble: list[Partition] = gen_ensemble(chain, n_elections)
    for map in ensemble:
        with CodeTimer("running district elections"):
            election_results = run_statewide_election(map, mmd_config) 
        elections_results.append(election_results)
    return elections_results


def plot_party_split(elections_results: list[list[Candidate]]): 
    dem_counts: list[int] = []
    for election_result in elections_results:
        dem_count = 0
        for candidate in election_result:
            if candidate.party == Party.DEMOCRAT:
                dem_count += 1
        dem_counts.append(dem_count)
    print("DEM COUNTS: " + str(dem_counts))
    hist(dem_counts)
    plt.savefig('dem_split2.png')


def run_smd_election(partition: Partition):
    pass


def main() -> None:
    prs = Presentation() # perhaps make this a static variable inside plotting functions

    prec_graph: Graph = Graph.from_json(GRAPH_FILE)
    prec_geometries: GeoSeries = GeoSeries.from_file(SHAPE_FILE)

    smd_partition: Partition = Partition(prec_graph, assignment=SMD_ASSIGNMENT_COL)
    # mmd_config = random.choice(gen_mmd_configs(len(smd_partition.parts)))
    mmd_config = gen_mmd_configs(len(smd_partition.parts))[0]
    mmd_assignment = gen_mmd_seed_assignment(smd_partition, mmd_config)

    mmd_partition: Partition = Partition(
        graph=prec_graph, 
        assignment=mmd_assignment,
        updaters={"cut_edges": cut_edges, "population": Tally(POP_COL, "population")},
    )

    chain = MarkovChain(
        partial(mmd_recom, mmd_config=mmd_config, epsilon=0.01),
        [],
        always_accept,
        mmd_partition,
        total_steps=3
    )

    mmd_election_results: list[Candidate] = run_mmd_election(mmd_partition, 1, mmd_config)
    print(mmd_election_results)
    # ensemble: list[Partition] = gen_ensemble(chain, 10)
    # for map in ensemble:
    #     plot_partition(map, prec_geometries, prs, mmd_config, show=False, cmap=COLORS)
    # prs.save(PRES_FILE)
    # elections_results = run_many_statewide_elections(chain, mmd_config, 20)    
    # plot_party_split(elections_results)


if __name__ == "__main__":
    main()