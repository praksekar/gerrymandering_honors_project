from gerrychain import Graph, Partition, MarkovChain, constraints
from gerrychain.updaters import Tally
from gerrychain.constraints import single_flip_contiguous
from gerrychain.proposals import propose_random_flip, recom
from gerrychain.accept import always_accept
from functools import partial
import matplotlib.pyplot as plt
from pprint import pprint
import numpy as np
from pptx import Presentation


def add_plot_to_pres(xvals, yvals, colors, step_no, prs):
    scat = plt.scatter(xvals, yvals, s=4, c=colors)
    plt.title("Markov chain step: %d" % step_no)
    plt.savefig("/tmp/graph.png", dpi=100)
    scat.remove()
    slide = prs.slides.add_slide(blank_slide_layout)
    slide.shapes.add_picture("/tmp/graph.png", 0, 0)


def get_district_colors(curr_partition):
    num_nodes = len(curr_partition.graph.nodes)
    colors = np.array(['#e6194b', '#3cb44b', '#ffe119', '#4363d8', '#f58231', '#911eb4', '#46f0f0', '#f032e6', '#bcf60c', '#fabebe', '#008080',
                       '#e6beff', '#9a6324', '#fffac8', '#800000', '#aaffc3', '#808000', '#ffd8b1', '#000075', '#808080', '#ffffff', '#000000'])
    district_colors = [None]*num_nodes
    count = 0
    for id, node in curr_partition.graph.nodes(data=True):
        district_colors[count] = colors[curr_partition.assignment[id]]
        count += 1
    return district_colors


def get_node_coordinates(graph):
    num_nodes = len(graph.nodes)
    xvals = np.zeros(num_nodes)
    yvals = np.zeros(num_nodes)
    count = 0
    for id, node in graph.nodes(data=True):
        xvals[count] = float(node["INTPTLON10"])
        yvals[count] = float(node["INTPTLAT10"])
        count += 1
    return xvals, yvals


# main code ----------
data_dir = "./state_data/"
data_file = "PA_VTDs.json"
pres_dir = "./pres_output/"
pres_file = "slideshow4.pptx"
markov_chain_steps = 10000
save_interval = 1

plt.figure(figsize=(13.0, 8.0))  # plot config
plt.xlabel("Longitude(degrees)")
plt.ylabel("Latitude(degrees)")

graph = Graph.from_json("%s%s" % (data_dir, data_file))

seed_partition = Partition(
    graph,
    assignment="CD_2011",
    updaters={
        "population": Tally("TOTPOP", alias="population"),
    }
)

ideal_population = sum(
    seed_partition["population"].values()) / len(seed_partition)

proposal = partial(recom,
                   pop_col="TOTPOP",
                   pop_target=ideal_population,
                   epsilon=0.02,
                   node_repeats=2
                   )

compactness_bound = constraints.UpperBound(
    lambda p: len(p["cut_edges"]),
    2*len(seed_partition["cut_edges"])
)

pop_constraint = constraints.within_percent_of_ideal_population(
    seed_partition, 0.02)

chain = MarkovChain(
    proposal=proposal,
    constraints=[
        pop_constraint,
        compactness_bound
    ],
    accept=always_accept,
    initial_state=seed_partition,
    total_steps=markov_chain_steps
)

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
xvals, yvals = get_node_coordinates(seed_partition.graph)

try:
    count = 0
    for partition in chain:
        print("%s%d" % ("Markov chain step: ", count))
        district_colors = get_district_colors(curr_partition=partition)
        if (count % save_interval == 0):
            print("Added slide")
            add_plot_to_pres(xvals, yvals, district_colors, count, prs)
        count += 1
except KeyboardInterrupt:
    prs.save("%s%s" % (pres_dir, pres_file))
