from matplotlib.pyplot import hist
from geopandas import GeoSeries
from gerrychain import Partition 
import matplotlib.pyplot as plt
from ..custom_types import ElectionsResults
from .election import Candidate, Party
from pptx import Presentation
from ..custom_types import Ensemble
from .utils import is_path_in_proj
import logging 
logger = logging.getLogger(__name__)
import consts
from pathlib import Path


prs = Presentation() 


def add_plot_to_pres(prs) -> None:
    plt.savefig("/tmp/graph.png")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture("/tmp/graph.png", 0, 0)
    print("slide added")


def plot_partition(partition: Partition, prs: Presentation=None, show: bool = False) -> None:
    logger.info(f"plotting {partition}")
    partition.plot(cmap=consts.DISTINCT_COLORS)
    centroids: dict[int, tuple] = get_district_centroids(partition, partition.graph.geometry)
    for districtID, coord in centroids.items():
        pop_frac = float(partition[consts.POP_UPDATER][districtID]/sum(partition[consts.POP_UPDATER].values())) * sum(partition.district_reps.values())
        plt.text(coord[0], coord[1], "District %d\nPopulation: %d\nPop Frac: %3f/18\nnum reps: %d" % (districtID, partition[consts.POP_UPDATER][districtID], pop_frac, partition.district_reps[districtID]))
    if prs is not None:
        add_plot_to_pres(prs)
    if show:
        plt.show()


def plot_ensemble(ensemble: Ensemble, prs: Presentation=None, show: bool = False) -> None:
    for map in ensemble.maps:
        plot_partition(map, prs=prs, show=show)


def get_district_centroids(partition: Partition, precinct_geometries: GeoSeries) -> dict[int, tuple]:
    district_centroids = {}
    for part, nodes in partition.parts.items():
        part_geometries = precinct_geometries.loc[precinct_geometries.index[list(nodes)]]
        district_centroids[part] = tuple(part_geometries.unary_union.centroid.coords)[0]
    return district_centroids


def plot_party_split(elections_results: ElectionsResults, n_districts: int, file: Path): 
    dem_counts: list[int] = []
    for election_result in elections_results.results:
        dem_count = 0
        for candidate in election_result:
            if candidate.party == Party.DEMOCRAT:
                dem_count += 1
        dem_counts.append(dem_count)
    print("DEM COUNTS: " + str(dem_counts))
    hist(dem_counts, range=(0, n_districts))
    if not is_path_in_proj(file):
        raise Exception("attempting to write in file outside of project directory")
    logger.info(f"saving plot to {file}")
    file.parent.mkdir(exist_ok=True, parents=True)
    plt.savefig(file)


def vote_seat_share_curve():
    pass


def box_and_whisker_plot():
    pass